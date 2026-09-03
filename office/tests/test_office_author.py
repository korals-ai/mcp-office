"""Tests for the office authoring logic.

xlsx/docx/pptx authoring is pure-python and runs without LibreOffice, so it's
fully covered here (produce a file, verify its magic bytes, reopen it with the
real reader). author_pdf needs ``soffice`` (absent in CI) so it's skipped
there — the render fork is exercised by the post-deploy integration suite.
"""

from __future__ import annotations

import shutil
from pathlib import Path

import pytest

from src.office_author import (
    OfficeAuthorError,
    author_docx,
    author_pdf,
    author_pptx,
    author_xlsx,
)

_ZIP_MAGIC = b"PK\x03\x04"
_PDF_MAGIC = b"%PDF"


def _magic(path: Path, n: int = 4) -> bytes:
    with path.open("rb") as fh:
        return fh.read(n)


# --- xlsx -------------------------------------------------------------------


def test_author_xlsx_produces_a_valid_zip_and_reopens(tmp_path: Path) -> None:
    dest = tmp_path / "book.xlsx"
    out = author_xlsx(dest, [["name", "value"], ["apple", 1], ["banana", 2]])
    assert out == dest
    assert _magic(dest) == _ZIP_MAGIC  # a real OOXML package, not text

    from openpyxl import load_workbook

    ws = load_workbook(dest).active
    assert [c.value for c in ws[1]] == ["name", "value"]
    assert ws["B3"].value == 2
    assert ws[1][0].font.bold is True  # header bolded by default


def test_author_xlsx_creates_parent_dirs(tmp_path: Path) -> None:
    dest = tmp_path / "nested" / "deep" / "book.xlsx"
    author_xlsx(dest, [["a"]])
    assert dest.exists()


def test_author_xlsx_sanitizes_and_truncates_sheet_name(tmp_path: Path) -> None:
    from openpyxl import load_workbook

    dest = tmp_path / "s.xlsx"
    author_xlsx(dest, [["a"]], sheet_name="a/b:c" + "x" * 40)
    title = load_workbook(dest).active.title
    assert len(title) <= 31
    assert not set(title) & set(":\\/?*[]")


def test_author_xlsx_empty_rows_rejected(tmp_path: Path) -> None:
    with pytest.raises(OfficeAuthorError, match="empty"):
        author_xlsx(tmp_path / "x.xlsx", [])


# --- docx -------------------------------------------------------------------


def test_author_docx_produces_a_valid_zip_and_reopens(tmp_path: Path) -> None:
    dest = tmp_path / "doc.docx"
    author_docx(dest, "Report", ["First para.", "Second para."])
    assert _magic(dest) == _ZIP_MAGIC

    from docx import Document

    text = "\n".join(p.text for p in Document(str(dest)).paragraphs)
    assert "Report" in text and "First para." in text


def test_author_docx_requires_content(tmp_path: Path) -> None:
    with pytest.raises(OfficeAuthorError, match="title and/or"):
        author_docx(tmp_path / "d.docx", None, [])


# --- pptx -------------------------------------------------------------------


def test_author_pptx_produces_a_valid_zip_and_reopens(tmp_path: Path) -> None:
    dest = tmp_path / "deck.pptx"
    author_pptx(dest, [{"title": "Intro", "bullets": ["one", "two"]}, {"title": "End"}])
    assert _magic(dest) == _ZIP_MAGIC

    from pptx import Presentation

    prs = Presentation(str(dest))
    assert len(prs.slides) == 2
    assert prs.slides[0].shapes.title.text == "Intro"


def test_author_pptx_empty_slides_rejected(tmp_path: Path) -> None:
    with pytest.raises(OfficeAuthorError, match="empty"):
        author_pptx(tmp_path / "p.pptx", [])


# --- pdf (needs soffice) ----------------------------------------------------


@pytest.mark.skipif(shutil.which("soffice") is None, reason="LibreOffice not installed")
def test_author_pdf_produces_a_valid_pdf(tmp_path: Path) -> None:
    dest = tmp_path / "out.pdf"
    author_pdf(dest, "Title", ["Body paragraph."])
    assert _magic(dest) == _PDF_MAGIC
    # The intermediate DOCX must not be left behind next to the output.
    assert list(tmp_path.glob("*.docx")) == []
