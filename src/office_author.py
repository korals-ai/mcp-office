"""Author Office documents from structured data.

The sanctioned way for the workspace agent to CREATE ``.xlsx`` / ``.docx`` /
``.pptx`` / ``.pdf`` — the companion to :mod:`src.office_convert` (which only
*converts* an existing file). A valid file is produced **by construction**: the
agent hands over structured data (rows, paragraphs, slides) and a real writer
(openpyxl / python-docx / python-pptx) serialises the bytes, so the agent can
never emit the broken "XML-under-a-.xlsx-name" that the workspace write guard
exists to catch. PDF composes: author a DOCX, then reuse LibreOffice
``convert`` — no separate PDF engine.

Zero-copy data plane, same as ``convert``: the caller names a path on the shared
tenant PVC; this writes the file there and returns the path.
"""

from __future__ import annotations

import re
import shutil
import tempfile
from pathlib import Path
from typing import Any

from src.office_convert import OfficeConvertError
from src.office_convert import convert as _convert


class OfficeAuthorError(Exception):
    """Authoring failed — an empty/invalid spec, an unwritable path, or a
    library error. Surfaced to the agent as an MCP tool error."""


# Excel forbids these in a sheet title, and caps it at 31 chars.
_XLSX_TITLE_BAD = re.compile(r"[:\\/?*\[\]]")


def _ensure_parent(dest: Path) -> None:
    try:
        dest.parent.mkdir(parents=True, exist_ok=True)
    except OSError as exc:  # unwritable mount, path is a file, ...
        raise OfficeAuthorError(f"cannot create directory {dest.parent}: {exc}") from exc


def author_xlsx(
    dest: Path,
    rows: list[list[Any]],
    sheet_name: str = "Sheet1",
    bold_header: bool = True,
) -> Path:
    """Write ``rows`` (a list of rows, each a list of cell values) to a real
    ``.xlsx`` workbook at ``dest``. First row is bolded when ``bold_header``."""
    if not rows:
        raise OfficeAuthorError("author_xlsx: 'rows' is empty — provide at least one row.")

    from openpyxl import Workbook
    from openpyxl.styles import Font

    title = _XLSX_TITLE_BAD.sub("_", sheet_name).strip()[:31] or "Sheet1"
    wb = Workbook()
    ws = wb.active
    ws.title = title
    for row in rows:
        ws.append(list(row))
    if bold_header:
        for cell in ws[1]:
            cell.font = Font(bold=True)

    _ensure_parent(dest)
    try:
        wb.save(dest)
    except OSError as exc:
        raise OfficeAuthorError(f"cannot write {dest}: {exc}") from exc
    return dest


def author_docx(dest: Path, title: str | None, paragraphs: list[str]) -> Path:
    """Write a real ``.docx`` at ``dest``: an optional level-1 heading
    (``title``) followed by one body paragraph per entry in ``paragraphs``."""
    if not title and not paragraphs:
        raise OfficeAuthorError("author_docx: provide a title and/or at least one paragraph.")

    from docx import Document

    doc = Document()
    if title:
        doc.add_heading(title, level=1)
    for para in paragraphs:
        doc.add_paragraph(para)

    _ensure_parent(dest)
    try:
        doc.save(str(dest))
    except OSError as exc:
        raise OfficeAuthorError(f"cannot write {dest}: {exc}") from exc
    return dest


def author_pptx(dest: Path, slides: list[dict[str, Any]]) -> Path:
    """Write a real ``.pptx`` at ``dest``. ``slides`` is a list of
    ``{"title": str, "bullets": [str, ...]}`` — one Title-and-Content slide
    each; ``bullets`` is optional."""
    if not slides:
        raise OfficeAuthorError("author_pptx: 'slides' is empty — provide at least one slide.")

    from pptx import Presentation

    prs = Presentation()
    content_layout = prs.slide_layouts[1]  # "Title and Content"
    for spec in slides:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = str(spec.get("title", ""))
        bullets = spec.get("bullets") or []
        if bullets:
            frame = slide.placeholders[1].text_frame
            frame.text = str(bullets[0])
            for bullet in bullets[1:]:
                frame.add_paragraph().text = str(bullet)

    _ensure_parent(dest)
    try:
        prs.save(str(dest))
    except OSError as exc:
        raise OfficeAuthorError(f"cannot write {dest}: {exc}") from exc
    return dest


def author_pdf(dest: Path, title: str | None, paragraphs: list[str]) -> Path:
    """Write a real ``.pdf`` at ``dest`` by authoring a DOCX and rendering it
    with LibreOffice — no separate PDF engine. Same content model as
    :func:`author_docx`."""
    _ensure_parent(dest)
    # Author + render inside a temp dir on the SAME filesystem as dest so the
    # final move is an atomic rename, and the intermediate DOCX never lands
    # next to the user's files.
    try:
        with tempfile.TemporaryDirectory(dir=dest.parent) as td:
            tmp = Path(td)
            author_docx(tmp / "doc.docx", title, paragraphs)
            try:
                rendered = _convert(tmp / "doc.docx", tmp, to="pdf")
            except OfficeConvertError as exc:
                detail = f": {exc.stderr.strip()}" if exc.stderr else ""
                raise OfficeAuthorError(f"PDF render failed{detail}") from exc
            shutil.move(str(rendered), str(dest))
    except OSError as exc:
        raise OfficeAuthorError(f"cannot write {dest}: {exc}") from exc
    return dest
